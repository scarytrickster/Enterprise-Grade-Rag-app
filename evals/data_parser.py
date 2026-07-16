"""
Parses true_data (all files) and noisy_data (pptx/docx/txt only) into tagged chunks.
Uses python-docx and python-pptx directly — bypasses unstructured to avoid segfaults.
Reuses parse_text, parse_html, and chunk_text from the main app.
"""

import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import docx as python_docx
from pptx import Presentation

from app.ingestion.loaders.text import parse_text
from app.ingestion.loaders.html import parse_html
from app.ingestion.chunking.splitter import chunk_text

TRUE_DATA_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "DATA", "true_data")
NOISY_DATA_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "DATA", "noisy_data")
NOISY_ALLOWED_EXTS = {".pptx", ".docx", ".txt"}


def _parse_docx(file_path: str) -> str:
    doc = python_docx.Document(file_path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
    return "\n".join(texts)


def parse_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".docx":
            return _parse_docx(file_path)
        elif ext == ".pptx":
            return _parse_pptx(file_path)
        elif ext in (".txt", ".md"):
            return parse_text(file_path)
        elif ext in (".html", ".htm"):
            return parse_html(file_path)
    except Exception:
        pass
    return ""


def load_all_chunks() -> list[dict]:
    """
    Returns all chunks tagged with source filename and whether they are noise.
    Used by the eval pipeline to understand what context the RAG system draws from.
    """
    results = []

    for fname in sorted(os.listdir(TRUE_DATA_DIR)):
        fpath = os.path.join(TRUE_DATA_DIR, fname)
        if not os.path.isfile(fpath):
            continue
        text = parse_file(fpath)
        if text:
            for chunk in chunk_text(text):
                results.append({"text": chunk, "source": fname, "is_noise": False})

    for fname in sorted(os.listdir(NOISY_DATA_DIR)):
        ext = os.path.splitext(fname)[1].lower()
        if ext not in NOISY_ALLOWED_EXTS:
            continue
        fpath = os.path.join(NOISY_DATA_DIR, fname)
        if not os.path.isfile(fpath):
            continue
        text = parse_file(fpath)
        if text:
            for chunk in chunk_text(text):
                results.append({"text": chunk, "source": fname, "is_noise": True})

    return results
